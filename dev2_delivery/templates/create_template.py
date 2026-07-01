"""
Run once to generate the base template.pptx with slide layouts and branding.
After running, template.pptx lives in dev2_delivery/templates/.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT = RGBColor(0x7C, 0x3A, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_accent_bar(slide):
    """Purple bar on the left edge for branding."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        Inches(0.15), SLIDE_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank = prs.slide_layouts[6]  # completely blank layout

# --- Slide 1: Title ---
s1 = prs.slides.add_slide(blank)
set_bg(s1, DARK_BG)
add_accent_bar(s1)
add_textbox(s1, "COMPANY RESEARCH REPORT",
            Inches(0.4), Inches(2.5), Inches(8), Inches(1),
            font_size=14, color=ACCENT)
add_textbox(s1, "{{company_name}}",
            Inches(0.4), Inches(3.1), Inches(10), Inches(1.5),
            font_size=44, bold=True)
add_textbox(s1, "{{website_url}}",
            Inches(0.4), Inches(4.8), Inches(8), Inches(0.5),
            font_size=14, color=LIGHT_GRAY)

# --- Slide 2: About ---
s2 = prs.slides.add_slide(blank)
set_bg(s2, DARK_BG)
add_accent_bar(s2)
add_textbox(s2, "ABOUT", Inches(0.4), Inches(0.3),
            Inches(4), Inches(0.6), font_size=12, color=ACCENT)
add_textbox(s2, "Company Overview",
            Inches(0.4), Inches(0.9), Inches(12), Inches(0.8),
            font_size=28, bold=True)
add_textbox(s2, "{{about_summary}}",
            Inches(0.4), Inches(1.9), Inches(12), Inches(2),
            font_size=16, color=LIGHT_GRAY)
add_textbox(s2, "Industry: {{industry}}   |   Founded: {{founded}}   |   Size: {{size}}",
            Inches(0.4), Inches(6.2), Inches(12), Inches(0.6),
            font_size=13, color=ACCENT)

# --- Slide 3: Products & Services ---
s3 = prs.slides.add_slide(blank)
set_bg(s3, DARK_BG)
add_accent_bar(s3)
add_textbox(s3, "OFFERINGS", Inches(0.4), Inches(0.3),
            Inches(4), Inches(0.6), font_size=12, color=ACCENT)
add_textbox(s3, "Products & Services",
            Inches(0.4), Inches(0.9), Inches(12), Inches(0.8),
            font_size=28, bold=True)
add_textbox(s3, "Products\n{{products}}",
            Inches(0.4), Inches(2.0), Inches(5.5), Inches(4),
            font_size=15, color=LIGHT_GRAY)
add_textbox(s3, "Services\n{{services}}",
            Inches(7.0), Inches(2.0), Inches(5.5), Inches(4),
            font_size=15, color=LIGHT_GRAY)

# --- Slide 4: News ---
s4 = prs.slides.add_slide(blank)
set_bg(s4, DARK_BG)
add_accent_bar(s4)
add_textbox(s4, "LATEST NEWS", Inches(0.4), Inches(0.3),
            Inches(4), Inches(0.6), font_size=12, color=ACCENT)
add_textbox(s4, "Recent Developments",
            Inches(0.4), Inches(0.9), Inches(12), Inches(0.8),
            font_size=28, bold=True)
add_textbox(s4, "{{news}}",
            Inches(0.4), Inches(2.0), Inches(12), Inches(4.5),
            font_size=15, color=LIGHT_GRAY)

# --- Slide 5: LLM Analysis ---
s5 = prs.slides.add_slide(blank)
set_bg(s5, DARK_BG)
add_accent_bar(s5)
add_textbox(s5, "ANALYSIS", Inches(0.4), Inches(0.3),
            Inches(4), Inches(0.6), font_size=12, color=ACCENT)
add_textbox(s5, "AI-Powered Insights",
            Inches(0.4), Inches(0.9), Inches(12), Inches(0.8),
            font_size=28, bold=True)
add_textbox(s5, "Pain Points\n{{pain_points}}",
            Inches(0.4), Inches(2.0), Inches(3.8), Inches(4),
            font_size=14, color=LIGHT_GRAY)
add_textbox(s5, "Growth Signals\n{{growth_signals}}",
            Inches(4.7), Inches(2.0), Inches(3.8), Inches(4),
            font_size=14, color=LIGHT_GRAY)
add_textbox(s5, "Tech Stack\n{{tech_stack}}",
            Inches(9.0), Inches(2.0), Inches(3.8), Inches(4),
            font_size=14, color=LIGHT_GRAY)

# --- Slide 6: Recommended Services ---
s6 = prs.slides.add_slide(blank)
set_bg(s6, DARK_BG)
add_accent_bar(s6)
add_textbox(s6, "OUR RECOMMENDATIONS", Inches(0.4), Inches(0.3),
            Inches(6), Inches(0.6), font_size=12, color=ACCENT)
add_textbox(s6, "How We Can Help",
            Inches(0.4), Inches(0.9), Inches(12), Inches(0.8),
            font_size=28, bold=True)
add_textbox(s6, "{{recommended_services}}",
            Inches(0.4), Inches(2.0), Inches(12), Inches(4.5),
            font_size=15, color=LIGHT_GRAY)

import os
out_path = os.path.join(os.path.dirname(__file__), "template.pptx")
prs.save(out_path)
print(f"Template saved to {out_path}")