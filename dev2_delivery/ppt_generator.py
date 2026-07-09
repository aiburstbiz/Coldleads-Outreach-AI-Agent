"""
ppt_generator.py — Dev2 service
Consumes a CompanyResearch object and produces a filled .pptx file.

Template: AIBurst-PitchDeck-2.pptx (9 slides)
  1. Title (dynamic: Company Name, Company Logo)
  2. About (dynamic)
  3. Founder Profile (STATIC — same on every deck)
  4. Offerings / Products & Services (dynamic)
  5. Latest News (dynamic)
  6. Analysis (dynamic)
  7. Recommendations (dynamic)
  8. Our Team (STATIC — same on every deck)
  9. Closing / Contact (STATIC — same on every deck)

Notes:
  - {{Subheading}} placeholder is removed entirely (dropped per Day 11 decision).
  - Company Logo is downloaded from data.logo_url (validated Hunter.io URL from
    dev1_research/logo.py) and inserted as an image, preserving aspect ratio.
    If logo_url is None or the download fails, the placeholder is removed and
    the slide is left without a logo rather than erroring out.
  - Day 8's snapshot/spotlight slides are NOT used with this template (dropped
    per Day 11 decision) — this template has no corresponding slide for them.
"""
import io
import os
import uuid

import requests
from pptx import Presentation
from pptx.util import Inches
from shared.schema import CompanyResearch

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "templates", "template.pptx")
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "outputs")

LOGO_DOWNLOAD_TIMEOUT = 5


# ── text helpers ──────────────────────────────────────────────────────────────

def _replace_text(shape, replacements: dict):
    """Replace {{key}} placeholders in a shape's text frame."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for key, val in replacements.items():
                if key in run.text:
                    run.text = run.text.replace(key, val)


def _format_list(items: list, prefix="• ") -> str:
    if not items:
        return "None found"
    return "\n".join(f"{prefix}{item}" for item in items)


def _format_news(news_items: list) -> str:
    if not news_items:
        return "No recent news found"
    lines = []
    for item in news_items:
        date = f" ({item.date})" if item.date else ""
        lines.append(f"• {item.title}{date}\n  {item.summary}")
    return "\n\n".join(lines)


def _format_recommended_services(services: list) -> str:
    if not services:
        return "No recommendations generated"
    lines = []
    for svc in services:
        priority_label = f"[{svc.priority.value.upper()}]"
        lines.append(f"• {svc.service} {priority_label}\n  {svc.reason}")
    return "\n\n".join(lines)


# ── shape helpers ─────────────────────────────────────────────────────────────

def _find_shapes_with_text(slide, token: str) -> list:
    """Return all shapes on a slide whose text frame contains the given token."""
    matches = []
    for shape in slide.shapes:
        if shape.has_text_frame and token in shape.text_frame.text:
            matches.append(shape)
    return matches


def _remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def _remove_placeholder_shapes(prs: Presentation, token: str) -> None:
    """Remove every shape (across all slides) whose text is exactly/contains
    the given placeholder token. Used for placeholders we've decided to drop
    entirely (e.g. {{Subheading}})."""
    for slide in prs.slides:
        for shape in _find_shapes_with_text(slide, token):
            _remove_shape(shape)


LOGO_TARGET_HEIGHT = Inches(0.65)  # matches the AIBurst brand logo size used elsewhere in the deck


def _insert_logo(prs: Presentation, logo_url: str | None) -> None:
    """
    Finds the {{Company Logo}} placeholder shape, captures its center point,
    removes the placeholder, and — if logo_url is available — downloads the
    image and inserts it centered on that same point. Sized to
    LOGO_TARGET_HEIGHT (matching the AIBurst brand logo's size elsewhere in
    the deck) rather than the original placeholder's small box; width is
    auto-computed by python-pptx to preserve the image's aspect ratio.

    If logo_url is None or the download fails for any reason, the slide is
    simply left without a logo (placeholder removed, no image inserted) —
    this must never raise or block deck generation.
    """
    token = "{{Company Logo}}"
    for slide in prs.slides:
        for shape in _find_shapes_with_text(slide, token):
            center_x = shape.left + shape.width // 2
            center_y = shape.top + shape.height // 2
            _remove_shape(shape)

            if not logo_url:
                continue

            try:
                resp = requests.get(logo_url, timeout=LOGO_DOWNLOAD_TIMEOUT)
                if resp.status_code == 200 and resp.content:
                    image_stream = io.BytesIO(resp.content)
                    pic = slide.shapes.add_picture(image_stream, 0, 0, height=LOGO_TARGET_HEIGHT)
                    pic.left = center_x - pic.width // 2
                    pic.top = center_y - pic.height // 2
            except Exception:
                # Logo is a nice-to-have, not critical — never let a network
                # hiccup here break PPT generation.
                pass


# ── main entry point ──────────────────────────────────────────────────────────

def generate_pptx(data: CompanyResearch) -> str:
    """
    Generate a filled .pptx from a CompanyResearch object.
    Returns the absolute path to the saved file.
    """
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    prs = Presentation(TEMPLATE_PATH)

    replacements = {
        "{title}": f"AI-Integration Pipeline for {data.company_name}",
        "{{Company Name}}": data.company_name,
        "{{about_summary}}": data.about.summary if data.about else "N/A",
        "{{industry}}": data.about.industry if data.about else "N/A",
        "{{founded}}": data.about.founded or "N/A",
        "{{size}}": data.about.size or "N/A",
        "{{products}}": _format_list(data.products),
        "{{services}}": _format_list(data.services),
        "{{news}}": _format_news(data.news),
        "{{pain_points}}": _format_list(data.llm_analysis.pain_points),
        "{{growth_signals}}": _format_list(data.llm_analysis.growth_signals),
        "{{tech_stack}}": _format_list(data.llm_analysis.tech_stack_hints),
        "{{recommended_services}}": _format_recommended_services(data.recommended_services),
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_text(shape, replacements)

    # Dropped placeholders / features (Day 11 decisions)
    _remove_placeholder_shapes(prs, "{{Subheading}}")
    _insert_logo(prs, getattr(data, "logo_url", None))

    filename = f"{data.company_name.replace(' ', '_')}_{uuid.uuid4().hex[:8]}.pptx"
    output_path = os.path.join(OUTPUT_DIR, filename)
    prs.save(output_path)
    return output_path
